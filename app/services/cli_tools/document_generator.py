"""
Document Generation Service
Implements document creation using Pandoc, LaTeX, Markdown, and Python packages
"""

import asyncio
import tempfile
from typing import Dict, Any, List, Optional
from pathlib import Path
import json
import logging

from .interfaces import IDocumentGenerator, BaseGenerator, GenerationResult

logger = logging.getLogger(__name__)

class DocumentGenerationService(BaseGenerator, IDocumentGenerator):
    """Service for generating various document formats"""
    
    def __init__(self, output_dir: str = "/tmp/cli_generations/documents"):
        super().__init__(output_dir)
        self.required_tools = ["pandoc", "pdflatex", "xelatex"]
    
    async def generate_markdown(self, content: str, output_path: str) -> GenerationResult:
        """Generate markdown document"""
        try:
            if not output_path.endswith('.md'):
                output_path += '.md'
            
            full_path = self._get_output_path(output_path)
            
            with open(full_path, 'w', encoding='utf-8') as f:
                f.write(content)
            
            return self._create_generation_result(
                success=True,
                file_path=full_path,
                format="markdown",
                size=len(content)
            )
            
        except Exception as e:
            logger.error(f"Markdown generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )
    
    async def generate_pdf_from_markdown(self, markdown_content: str, output_path: str) -> GenerationResult:
        """Generate PDF from markdown using Pandoc"""
        try:
            if not output_path.endswith('.pdf'):
                output_path += '.pdf'
            
            full_path = self._get_output_path(output_path)
            
            # Create temporary markdown file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as temp_md:
                temp_md.write(markdown_content)
                temp_md_path = temp_md.name
            
            # Use Pandoc to convert markdown to PDF
            command = [
                "pandoc",
                temp_md_path,
                "-o", full_path,
                "--pdf-engine=xelatex",
                "--variable=geometry:margin=1in"
            ]
            
            success, stdout, stderr = await self._run_command(command)
            
            # Clean up temporary file
            Path(temp_md_path).unlink()
            
            if success and Path(full_path).exists():
                return self._create_generation_result(
                    success=True,
                    file_path=full_path,
                    format="pdf",
                    engine="pandoc"
                )
            else:
                return self._create_generation_result(
                    success=False,
                    error=f"Pandoc error: {stderr}"
                )
                
        except Exception as e:
            logger.error(f"PDF generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )
    
    async def generate_latex(self, content: str, output_path: str) -> GenerationResult:
        """Generate LaTeX document"""
        try:
            if not output_path.endswith('.tex'):
                output_path += '.tex'
            
            full_path = self._get_output_path(output_path)
            
            # Wrap content in basic LaTeX structure if not already wrapped
            if not content.strip().startswith("\\documentclass"):
                latex_content = f"""\\documentclass{{article}}
\\usepackage[utf8]{{inputenc}}
\\usepackage{{geometry}}
\\geometry{{margin=1in}}
\\title{{Generated Document}}
\\date{{\\today}}
\\begin{{document}}
\\maketitle

{content}

\\end{{document}}"""
            else:
                latex_content = content
            
            with open(full_path, 'w', encoding='utf-8') as f:
                f.write(latex_content)
            
            return self._create_generation_result(
                success=True,
                file_path=full_path,
                format="latex",
                size=len(latex_content)
            )
            
        except Exception as e:
            logger.error(f"LaTeX generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )
    
    async def generate_pdf_from_latex(self, latex_content: str, output_path: str) -> GenerationResult:
        """Generate PDF from LaTeX"""
        try:
            if not output_path.endswith('.pdf'):
                output_path += '.pdf'
            
            full_path = self._get_output_path(output_path)
            
            # Create temporary LaTeX file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.tex', delete=False, encoding='utf-8') as temp_tex:
                temp_tex.write(latex_content)
                temp_tex_path = temp_tex.name
            
            # Get directory for LaTeX compilation
            work_dir = Path(temp_tex_path).parent
            
            # Use pdflatex to compile
            command = [
                "pdflatex",
                "-interaction=nonstopmode",
                "-output-directory", str(work_dir),
                temp_tex_path
            ]
            
            success, stdout, stderr = await self._run_command(command)
            
            # Move generated PDF to desired location
            temp_pdf = work_dir / f"{Path(temp_tex_path).stem}.pdf"
            if success and temp_pdf.exists():
                import shutil
                shutil.move(str(temp_pdf), full_path)
                
                # Clean up temporary files
                Path(temp_tex_path).unlink()
                for ext in ['.aux', '.log', '.out']:
                    temp_file = work_dir / f"{Path(temp_tex_path).stem}{ext}"
                    if temp_file.exists():
                        temp_file.unlink()
                
                return self._create_generation_result(
                    success=True,
                    file_path=full_path,
                    format="pdf",
                    engine="pdflatex"
                )
            else:
                return self._create_generation_result(
                    success=False,
                    error=f"LaTeX compilation error: {stderr}"
                )
                
        except Exception as e:
            logger.error(f"LaTeX PDF generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )
    
    async def generate_word_document(self, content: str, output_path: str) -> GenerationResult:
        """Generate Word document using python-docx"""
        try:
            # Import here to handle missing dependency gracefully
            try:
                from docx import Document
                from docx.shared import Inches
            except ImportError:
                return self._create_generation_result(
                    success=False,
                    error="python-docx package not installed. Install with: pip install python-docx"
                )
            
            if not output_path.endswith('.docx'):
                output_path += '.docx'
            
            full_path = self._get_output_path(output_path)
            
            doc = Document()
            doc.add_heading('Generated Document', 0)
            
            # Split content into paragraphs
            paragraphs = content.split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    doc.add_paragraph(paragraph.strip())
            
            doc.save(full_path)
            
            return self._create_generation_result(
                success=True,
                file_path=full_path,
                format="docx",
                paragraphs=len(paragraphs)
            )
            
        except Exception as e:
            logger.error(f"Word document generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )
    
    async def generate_powerpoint(self, slides_data: List[Dict[str, Any]], output_path: str) -> GenerationResult:
        """Generate PowerPoint presentation using python-pptx"""
        try:
            # Import here to handle missing dependency gracefully
            try:
                from pptx import Presentation
                from pptx.util import Inches
            except ImportError:
                return self._create_generation_result(
                    success=False,
                    error="python-pptx package not installed. Install with: pip install python-pptx"
                )
            
            if not output_path.endswith('.pptx'):
                output_path += '.pptx'
            
            full_path = self._get_output_path(output_path)
            
            prs = Presentation()
            
            for slide_data in slides_data:
                # Add slide with title and content layout
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = slide_data.get('title', 'Slide Title')
                
                # Add content
                if 'content' in slide_data:
                    if isinstance(slide_data['content'], list):
                        tf = content.text_frame
                        tf.text = slide_data['content'][0] if slide_data['content'] else ''
                        
                        for item in slide_data['content'][1:]:
                            p = tf.add_paragraph()
                            p.text = item
                    else:
                        content.text = slide_data['content']
            
            prs.save(full_path)
            
            return self._create_generation_result(
                success=True,
                file_path=full_path,
                format="pptx",
                slides=len(slides_data)
            )
            
        except Exception as e:
            logger.error(f"PowerPoint generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )
    
    async def generate_excel(self, data: Dict[str, List[Dict[str, Any]]], output_path: str) -> GenerationResult:
        """Generate Excel spreadsheet using openpyxl"""
        try:
            # Import here to handle missing dependency gracefully
            try:
                from openpyxl import Workbook
                from openpyxl.utils import get_column_letter
            except ImportError:
                return self._create_generation_result(
                    success=False,
                    error="openpyxl package not installed. Install with: pip install openpyxl"
                )
            
            if not output_path.endswith('.xlsx'):
                output_path += '.xlsx'
            
            full_path = self._get_output_path(output_path)
            
            wb = Workbook()
            
            # Remove default sheet if we have data
            if data:
                wb.remove(wb.active)
            
            for sheet_name, sheet_data in data.items():
                ws = wb.create_sheet(title=sheet_name)
                
                if sheet_data:
                    # Add headers
                    headers = list(sheet_data[0].keys())
                    for col, header in enumerate(headers, 1):
                        ws[f"{get_column_letter(col)}1"] = header
                    
                    # Add data rows
                    for row, record in enumerate(sheet_data, 2):
                        for col, header in enumerate(headers, 1):
                            ws[f"{get_column_letter(col)}{row}"] = record.get(header, "")
            
            wb.save(full_path)
            
            return self._create_generation_result(
                success=True,
                file_path=full_path,
                format="xlsx",
                sheets=len(data),
                total_rows=sum(len(sheet_data) for sheet_data in data.values())
            )
            
        except Exception as e:
            logger.error(f"Excel generation error: {e}")
            return self._create_generation_result(
                success=False,
                error=str(e)
            )